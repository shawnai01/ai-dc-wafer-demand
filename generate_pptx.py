#!/usr/bin/env python3
"""AI DC 반도체 웨이퍼 수요 전망 — Executive PPTX Generator"""

import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.ticker as ticker
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from io import BytesIO
import os

# ─── Theme ───
BG = '#0F1117'
CARD_BG = '#1A1D27'
TEXT_W = '#FFFFFF'
TEXT_G = '#9CA3AF'
ACCENT_BLUE = '#3B82F6'
ACCENT_GREEN = '#10B981'
ACCENT_RED = '#EF4444'
ACCENT_YELLOW = '#F59E0B'
ACCENT_PURPLE = '#8B5CF6'
BULL_C = '#10B981'
BASE_C = '#3B82F6'
BEAR_C = '#EF4444'

plt.rcParams.update({
    'figure.facecolor': BG,
    'axes.facecolor': CARD_BG,
    'text.color': TEXT_W,
    'axes.labelcolor': TEXT_W,
    'xtick.color': TEXT_G,
    'ytick.color': TEXT_G,
    'axes.edgecolor': '#333',
    'grid.color': '#333',
    'grid.alpha': 0.5,
    'font.size': 11,
    'axes.unicode_minus': False,
})

# Try Korean font
for f in ['Apple SD Gothic Neo', 'AppleGothic', 'NanumGothic', 'Malgun Gothic', 'sans-serif']:
    try:
        plt.rcParams['font.family'] = f
        fig_test = plt.figure()
        plt.close(fig_test)
        break
    except:
        continue

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def rgb(hex_str):
    h = hex_str.lstrip('#')
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def add_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = rgb(BG)

def add_text(slide, left, top, width, height, text, size=14, bold=False, color=TEXT_W, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = rgb(color)
    p.alignment = align
    return txBox

def add_chart_image(slide, fig, left, top, width, height):
    buf = BytesIO()
    fig.savefig(buf, format='png', dpi=180, bbox_inches='tight', pad_inches=0.3)
    buf.seek(0)
    slide.shapes.add_picture(buf, Inches(left), Inches(top), Inches(width), Inches(height))
    plt.close(fig)

def add_shape_box(slide, left, top, width, height, fill_color=CARD_BG):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill_color)
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

# ══════════════════════════════════════════════
# SLIDE 1: Title
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 1, 1.5, 11, 1, 'AI 데이터센터', 48, True, ACCENT_BLUE, PP_ALIGN.CENTER)
add_text(slide, 1, 2.5, 11, 1, '반도체 웨이퍼 수요 전망', 42, True, TEXT_W, PP_ALIGN.CENTER)
add_text(slide, 1, 3.8, 11, 0.6, '2025–2027 시나리오별 수급 분석 및 투자 시사점', 20, False, TEXT_G, PP_ALIGN.CENTER)

# Key numbers
for i, (label, value) in enumerate([
    ('2025 반도체 매출', '$7,917억'),
    ('2026 전망', '$1조+'),
    ('빅테크 AI Capex', '$270–305B'),
]):
    x = 2.5 + i * 3
    add_shape_box(slide, x, 5.0, 2.5, 1.2)
    add_text(slide, x, 5.05, 2.5, 0.5, label, 11, False, TEXT_G, PP_ALIGN.CENTER)
    add_text(slide, x, 5.5, 2.5, 0.6, value, 22, True, ACCENT_BLUE, PP_ALIGN.CENTER)

add_text(slide, 1, 6.8, 11, 0.4, 'Semiconductor Research Team  |  2026.02.28  |  Confidential', 12, False, TEXT_G, PP_ALIGN.CENTER)

# ══════════════════════════════════════════════
# SLIDE 2: Scenario Overview — Grouped Bar Chart
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.5, 0.3, 12, 0.6, 'Bull / Base / Bear 시나리오 핵심 수치', 28, True)
add_text(slide, 0.5, 0.9, 12, 0.4, 'AI capex 성장률에 따라 2027년 웨이퍼 수요 2배+ 격차', 14, False, TEXT_G)

fig, axes = plt.subplots(1, 3, figsize=(12, 4))
categories = ['2025', '2027']
# AI Logic wafers
vals = {'Bull': [400, 750], 'Base': [330, 550], 'Bear': [270, 380]}
x = np.arange(2)
w = 0.25
for i, (label, v, c) in enumerate(zip(vals.keys(), vals.values(), [BULL_C, BASE_C, BEAR_C])):
    axes[0].bar(x + i*w, v, w, label=label, color=c, edgecolor='none')
axes[0].set_title('AI向 로직 웨이퍼 (천장/년)', fontsize=11, pad=8)
axes[0].set_xticks(x + w)
axes[0].set_xticklabels(categories)
axes[0].legend(fontsize=9)
axes[0].yaxis.set_major_formatter(ticker.FuncFormatter(lambda x,_: f'{int(x)}K'))

# HBM wafers
vals = {'Bull': [1200, 3000], 'Base': [1000, 2500], 'Bear': [750, 1500]}
for i, (label, v, c) in enumerate(zip(vals.keys(), vals.values(), [BULL_C, BASE_C, BEAR_C])):
    axes[1].bar(x + i*w, v, w, label=label, color=c, edgecolor='none')
axes[1].set_title('HBM 웨이퍼 (천장/년)', fontsize=11, pad=8)
axes[1].set_xticks(x + w)
axes[1].set_xticklabels(categories)
axes[1].legend(fontsize=9)
axes[1].yaxis.set_major_formatter(ticker.FuncFormatter(lambda x,_: f'{int(x)}K'))

# AI Semiconductor market
vals = {'Bull': [1800, 3500], 'Base': [1500, 2700], 'Bear': [1200, 1800]}
for i, (label, v, c) in enumerate(zip(vals.keys(), vals.values(), [BULL_C, BASE_C, BEAR_C])):
    axes[2].bar(x + i*w, v, w, label=label, color=c, edgecolor='none')
axes[2].set_title('AI 반도체 시장 ($억)', fontsize=11, pad=8)
axes[2].set_xticks(x + w)
axes[2].set_xticklabels(categories)
axes[2].legend(fontsize=9)

for ax in axes:
    ax.grid(axis='y', alpha=0.3)
    ax.set_axisbelow(True)
plt.tight_layout()
add_chart_image(slide, fig, 0.5, 1.4, 12.3, 5.5)

# ══════════════════════════════════════════════
# SLIDE 3: Bottleneck Priority
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.5, 0.3, 12, 0.6, '병목 우선순위 — CoWoS → HBM → 원재료', 28, True)
add_text(slide, 0.5, 0.9, 12, 0.4, 'HBM은 전 시나리오에서 2027까지 만성 부족', 14, False, TEXT_G)

fig, ax = plt.subplots(figsize=(12, 4.5))
bottlenecks = ['CoWoS', 'HBM', '첨단 노드\n(≤5nm)', '300mm\n원재료']
years = ['2025', '2026', '2027']
# Status: -1=부족, 0=보통, 1=여유
status = [
    [-1, 0, 1],   # CoWoS
    [-1, -1, -1],  # HBM
    [1, 1, 1],     # 첨단 노드
    [1, 1, 1],     # 원재료 (Base)
]
colors_map = {-1: ACCENT_RED, 0: ACCENT_YELLOW, 1: ACCENT_GREEN}
labels_map = {-1: '부족', 0: '보통', 1: '여유'}

for i, bn in enumerate(bottlenecks):
    for j, yr in enumerate(years):
        s = status[i][j]
        c = colors_map[s]
        ax.barh(i, 0.8, left=j, color=c, edgecolor=BG, linewidth=2, height=0.6)
        ax.text(j + 0.4, i, labels_map[s], ha='center', va='center', fontsize=12, fontweight='bold', color='white')

ax.set_yticks(range(len(bottlenecks)))
ax.set_yticklabels(bottlenecks, fontsize=13)
ax.set_xticks([0.4, 1.4, 2.4])
ax.set_xticklabels(years, fontsize=13)
ax.invert_yaxis()
ax.set_xlim(-0.1, 3)
ax.spines['top'].set_visible(False)
ax.spines['right'].set_visible(False)
ax.spines['bottom'].set_visible(False)
plt.tight_layout()
add_chart_image(slide, fig, 0.5, 1.4, 12.3, 5.5)

# ══════════════════════════════════════════════
# SLIDE 4: AI Chip Wafer Demand Breakdown (Donut + Bar)
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.5, 0.3, 12, 0.6, 'AI칩별 웨이퍼 소모량 (2025F)', 28, True)
add_text(slide, 0.5, 0.9, 12, 0.4, 'NVIDIA가 로직의 42% — 단일 기업 의존도 극대화', 14, False, TEXT_G)

fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(12, 4.5), gridspec_kw={'width_ratios': [1, 1.3]})

# Donut chart - Logic breakdown
sizes = [140, 80, 50, 60]
labels = ['NVIDIA GPU\n140K (42%)', 'AI ASIC\n80K (24%)', 'AMD GPU\n50K (15%)', 'AI서버 CPU\n60K (18%)']
colors = [ACCENT_GREEN, ACCENT_BLUE, ACCENT_PURPLE, ACCENT_YELLOW]
wedges, texts = ax1.pie(sizes, labels=labels, colors=colors, startangle=90, 
                         textprops={'color': TEXT_W, 'fontsize': 10},
                         wedgeprops=dict(width=0.4, edgecolor=BG))
ax1.set_title('로직 웨이퍼 330K장/년', fontsize=12, pad=10, color=TEXT_W)

# Bar chart - Logic vs HBM
categories = ['AI向 로직', 'HBM (DRAM)', '총합계']
values = [330, 1000, 1330]
bar_colors = [ACCENT_BLUE, ACCENT_GREEN, TEXT_G]
bars = ax2.barh(categories, values, color=bar_colors, height=0.5, edgecolor='none')
for bar, val in zip(bars, values):
    ax2.text(bar.get_width() + 20, bar.get_y() + bar.get_height()/2, f'{val:,}K', 
             va='center', fontsize=13, fontweight='bold', color=TEXT_W)
ax2.set_xlim(0, 1600)
ax2.set_title('로직 vs HBM (천장/년)', fontsize=12, pad=10, color=TEXT_W)
ax2.xaxis.set_major_formatter(ticker.FuncFormatter(lambda x,_: f'{int(x)}K'))
ax2.spines['top'].set_visible(False)
ax2.spines['right'].set_visible(False)
plt.tight_layout()
add_chart_image(slide, fig, 0.5, 1.4, 12.3, 5.5)

# ══════════════════════════════════════════════
# SLIDE 5: Sensitivity (Tornado chart)
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.5, 0.3, 12, 0.6, '감도 분석 — 핵심 변수별 웨이퍼 수요 스윙', 28, True)
add_text(slide, 0.5, 0.9, 12, 0.4, 'Capex 성장률이 최대 변수 (로직 18–22%, HBM 16–20% 스윙)', 14, False, TEXT_G)

fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(12, 4.5))

# Tornado - Logic
variables = ['Capex 성장률\n±10%p', 'TSMC 수율\n±5%p', 'CoWoS 증설\n±10K/월']
low = [-100, -22, 0]
high = [120, 18, 0]
y = range(len(variables))
ax1.barh(y, high, color=ACCENT_GREEN, height=0.5, label='상방')
ax1.barh(y, low, color=ACCENT_RED, height=0.5, label='하방')
ax1.set_yticks(y)
ax1.set_yticklabels(variables, fontsize=11)
ax1.set_title('AI向 로직 웨이퍼 Δ (천장/년)', fontsize=12, pad=10)
ax1.axvline(0, color='white', linewidth=0.5)
ax1.legend(fontsize=9)
for i, (l, h) in enumerate(zip(low, high)):
    if l != 0: ax1.text(l-8, i, f'{l}K', ha='right', va='center', fontsize=10, color=ACCENT_RED)
    if h != 0: ax1.text(h+8, i, f'+{h}K', ha='left', va='center', fontsize=10, color=ACCENT_GREEN)

# Tornado - HBM
variables2 = ['Capex 성장률\n±10%p', 'HBM 전환율\n±5%p']
low2 = [-500, -350]
high2 = [400, 450]
y2 = range(len(variables2))
ax2.barh(y2, high2, color=ACCENT_GREEN, height=0.5, label='상방')
ax2.barh(y2, low2, color=ACCENT_RED, height=0.5, label='하방')
ax2.set_yticks(y2)
ax2.set_yticklabels(variables2, fontsize=11)
ax2.set_title('HBM 웨이퍼 Δ (천장/년)', fontsize=12, pad=10)
ax2.axvline(0, color='white', linewidth=0.5)
ax2.legend(fontsize=9)
for i, (l, h) in enumerate(zip(low2, high2)):
    ax2.text(l-15, i, f'{l}K', ha='right', va='center', fontsize=10, color=ACCENT_RED)
    ax2.text(h+15, i, f'+{h}K', ha='left', va='center', fontsize=10, color=ACCENT_GREEN)

for ax in [ax1, ax2]:
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.grid(axis='x', alpha=0.3)
plt.tight_layout()
add_chart_image(slide, fig, 0.5, 1.4, 12.3, 5.5)

# ══════════════════════════════════════════════
# SLIDE 6: Capex → Semiconductor Flow
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.5, 0.3, 12, 0.6, '빅테크 AI Capex → 반도체 전환 경로', 28, True)

fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(12, 4.5))

# Capex by company
companies = ['Microsoft', 'Google', 'Amazon', 'Meta']
capex_25 = [63.5, 52.5, 67.5, 47.5]
capex_26_base = [65, 55, 75, 50]
capex_26_bull = [80, 70, 90, 60]

x = np.arange(len(companies))
w = 0.25
ax1.bar(x - w, capex_25, w, label='2025', color=TEXT_G, edgecolor='none')
ax1.bar(x, capex_26_base, w, label='2026 Base', color=BASE_C, edgecolor='none')
ax1.bar(x + w, capex_26_bull, w, label='2026 Bull', color=BULL_C, edgecolor='none')
ax1.set_xticks(x)
ax1.set_xticklabels(companies, fontsize=11)
ax1.set_title('빅테크 AI Capex ($B)', fontsize=12, pad=10)
ax1.legend(fontsize=9)
ax1.grid(axis='y', alpha=0.3)

# Funnel
funnel_labels = ['총 Capex\n$280B', '서버/네트워킹\n45%', '반도체\n65%', '최종\n$82B']
funnel_values = [280, 126, 82, 82]
funnel_widths = [1.0, 0.75, 0.55, 0.45]
colors_f = [TEXT_G, ACCENT_BLUE, ACCENT_GREEN, ACCENT_GREEN]
for i, (lbl, val, w_f, c) in enumerate(zip(funnel_labels, funnel_values, funnel_widths, colors_f)):
    ax2.barh(3-i, w_f, color=c, height=0.6, left=(1-w_f)/2, edgecolor=BG, linewidth=1)
    ax2.text(0.5, 3-i, f'{lbl}\n${val}B', ha='center', va='center', fontsize=11, fontweight='bold', color='white')
ax2.set_xlim(0, 1)
ax2.set_ylim(-0.5, 4)
ax2.axis('off')
ax2.set_title('Capex → 반도체 전환 (Base 2026)', fontsize=12, pad=10)
plt.tight_layout()
add_chart_image(slide, fig, 0.5, 1.4, 12.3, 5.5)

# ══════════════════════════════════════════════
# SLIDE 7: CoWoS Quarterly Trajectory
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.5, 0.3, 12, 0.6, 'CoWoS 분기별 수급 Trajectory (Base)', 28, True)
add_text(slide, 0.5, 0.9, 12, 0.4, '2026Q3 수급 균형 도달 — Bull에서는 2027까지 부족 지속', 14, False, TEXT_G)

fig, ax = plt.subplots(figsize=(12, 4.5))
quarters = ['25Q1','25Q2','25Q3','25Q4','26Q1','26Q2','26Q3','26Q4','27Q1','27Q2','27Q3','27Q4']
supply = [36,38,40,42,48,52,58,62,70,75,78,80]
demand = [42,44,45,46,50,54,56,58,65,70,75,78]

ax.plot(quarters, supply, '-o', color=ACCENT_GREEN, linewidth=2.5, markersize=6, label='공급')
ax.plot(quarters, demand, '-s', color=ACCENT_RED, linewidth=2.5, markersize=6, label='수요')
ax.fill_between(range(len(quarters)), supply, demand, where=[s<d for s,d in zip(supply,demand)], 
                alpha=0.15, color=ACCENT_RED, label='부족')
ax.fill_between(range(len(quarters)), supply, demand, where=[s>=d for s,d in zip(supply,demand)], 
                alpha=0.15, color=ACCENT_GREEN, label='여유')
ax.axvline(6, color=ACCENT_YELLOW, linestyle='--', alpha=0.7)
ax.text(6.1, max(demand)*0.95, '← 병목 해소\n   2026Q3', fontsize=10, color=ACCENT_YELLOW)
ax.set_ylabel('천 장/월', fontsize=12)
ax.legend(fontsize=10, loc='upper left')
ax.grid(axis='y', alpha=0.3)
ax.set_title('CoWoS 공급 vs 수요 (Base, 천 장/월)', fontsize=13, pad=10)
plt.tight_layout()
add_chart_image(slide, fig, 0.5, 1.4, 12.3, 5.5)

# ══════════════════════════════════════════════
# SLIDE 8: HBM Quarterly Trajectory
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.5, 0.3, 12, 0.6, 'HBM 분기별 수급 Trajectory (Base)', 28, True)
add_text(slide, 0.5, 0.9, 12, 0.4, '2027Q4까지 전 분기 부족 — 16-Hi 전환 시 비선형 확대', 14, False, TEXT_G)

fig, ax = plt.subplots(figsize=(12, 4.5))
hbm_supply = [240,250,260,280,330,370,410,450,490,520,550,580]
hbm_demand = [260,280,310,350,370,400,430,460,500,530,570,600]

ax.plot(quarters, hbm_supply, '-o', color=ACCENT_GREEN, linewidth=2.5, markersize=6, label='공급')
ax.plot(quarters, hbm_demand, '-s', color=ACCENT_RED, linewidth=2.5, markersize=6, label='수요')
ax.fill_between(range(len(quarters)), hbm_supply, hbm_demand, alpha=0.15, color=ACCENT_RED)
# Gap annotations
for i in [3, 7, 11]:
    gap = hbm_demand[i] - hbm_supply[i]
    mid = (hbm_supply[i] + hbm_demand[i]) / 2
    ax.annotate(f'-{gap}K', (i, mid), fontsize=9, color=ACCENT_RED, fontweight='bold', ha='center')

ax.set_ylabel('천 장/월', fontsize=12)
ax.legend(fontsize=10, loc='upper left')
ax.grid(axis='y', alpha=0.3)
ax.set_title('HBM DRAM 웨이퍼 공급 vs 수요 (Base, 천 장/월)', fontsize=13, pad=10)
plt.tight_layout()
add_chart_image(slide, fig, 0.5, 1.4, 12.3, 5.5)

# ══════════════════════════════════════════════
# SLIDE 9: Foundry Capa + Competition
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.5, 0.3, 12, 0.6, '파운드리 경쟁 — 점유율 재배분 시나리오', 28, True)
add_text(slide, 0.5, 0.9, 12, 0.4, 'Intel/삼성 성공해도 총 웨이퍼 수요 불변 — TSMC 최상위 노드 독점', 14, False, TEXT_G)

fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(12, 4.5))

# Stacked bar - Foundry capa 2025 vs 2027
labels_f = ['2025', '2027\nBase', '2027\nIntel성공', '2027\n복합성공']
tsmc = [340, 420, 390, 370]
samsung = [80, 110, 105, 130]
intel = [20, 40, 80, 80]
x_f = np.arange(len(labels_f))
ax1.bar(x_f, tsmc, 0.5, label='TSMC', color=ACCENT_BLUE, edgecolor='none')
ax1.bar(x_f, samsung, 0.5, bottom=tsmc, label='삼성', color=ACCENT_PURPLE, edgecolor='none')
ax1.bar(x_f, intel, 0.5, bottom=[t+s for t,s in zip(tsmc, samsung)], label='Intel', color=ACCENT_YELLOW, edgecolor='none')
ax1.set_xticks(x_f)
ax1.set_xticklabels(labels_f, fontsize=10)
ax1.set_title('첨단 노드 Capa (K/월)', fontsize=12, pad=10)
ax1.legend(fontsize=9)
ax1.grid(axis='y', alpha=0.3)

# Pie - Market share scenarios
scenarios = ['Base 2027', 'Intel 성공', '복합 성공']
shares = [[75, 18, 7], [69, 17, 14], [62.5, 24, 13.5]]
colors_pie = [ACCENT_BLUE, ACCENT_PURPLE, ACCENT_YELLOW]

for i, (scen, sh) in enumerate(zip(scenarios, shares)):
    ax_sub = fig.add_axes([0.55 + i*0.16, 0.15, 0.14, 0.65])
    ax_sub.set_facecolor(CARD_BG)
    wedges, _ = ax_sub.pie(sh, colors=colors_pie, startangle=90, 
                           wedgeprops=dict(width=0.4, edgecolor=BG))
    ax_sub.set_title(scen, fontsize=9, color=TEXT_W, pad=3)
    ax_sub.text(0, 0, f'TSMC\n{sh[0]}%', ha='center', va='center', fontsize=9, color=TEXT_W, fontweight='bold')

ax2.axis('off')
ax2.set_title('점유율 시나리오 (≤5nm)', fontsize=12, pad=10)
plt.tight_layout()
add_chart_image(slide, fig, 0.5, 1.4, 12.3, 5.5)

# ══════════════════════════════════════════════
# SLIDE 10: Risk Quantification
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.5, 0.3, 12, 0.6, '리스크 정량화 — 확률 × 영향', 28, True)

fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(12, 4.5))

# Bubble chart - individual risks
risks = ['AI capex\nROI 실망', '대만해협', '에너지 제약', 'HBM/N2\n수율실패', '중동 투자↑', 'AI Agent↑']
probs = [25, 7, 35, 20, 45, 35]
impacts = [-84, -58, -9, 5, 13, 14]
sizes_b = [600, 294, 23, 100, 90, 35]
colors_b = [ACCENT_RED, ACCENT_RED, ACCENT_YELLOW, ACCENT_YELLOW, ACCENT_GREEN, ACCENT_GREEN]

for r, p, imp, s, c in zip(risks, probs, impacts, sizes_b, colors_b):
    ax1.scatter(p, imp, s=abs(s)*3+100, color=c, alpha=0.7, edgecolors='white', linewidth=0.5)
    ax1.annotate(r, (p, imp), fontsize=8, color=TEXT_W, ha='center', va='bottom',
                xytext=(0, 10), textcoords='offset points')
ax1.axhline(0, color='white', linewidth=0.5, alpha=0.3)
ax1.set_xlabel('확률 (%)', fontsize=11)
ax1.set_ylabel('금액 영향 ($B)', fontsize=11)
ax1.set_title('개별 리스크 맵', fontsize=12, pad=10)
ax1.grid(alpha=0.3)

# Perfect Storm bars
storms = ['#1 AI버블\n+금융긴축', '#2 지정학\n+수율실패', '#3 수요과열\n+원재료병목']
storm_prob = [8, 3, 12]
storm_impact = [-150, -58, 8]
storm_colors = [ACCENT_RED, ACCENT_RED, ACCENT_YELLOW]
bars2 = ax2.barh(storms, storm_impact, color=storm_colors, height=0.5, edgecolor='none')
for bar, p in zip(bars2, storm_prob):
    x_pos = bar.get_width()
    offset = -10 if x_pos < 0 else 10
    ax2.text(x_pos + offset, bar.get_y() + bar.get_height()/2, 
             f'확률 {p}%', va='center', fontsize=10, color=TEXT_W, fontweight='bold',
             ha='right' if x_pos < 0 else 'left')
ax2.axvline(0, color='white', linewidth=0.5)
ax2.set_title('Perfect Storm 시나리오 ($B)', fontsize=12, pad=10)
ax2.set_xlabel('금액 영향 ($B)', fontsize=11)
ax2.grid(axis='x', alpha=0.3)
plt.tight_layout()
add_chart_image(slide, fig, 0.5, 1.4, 12.3, 5.5)

# ══════════════════════════════════════════════
# SLIDE 11: Investment — Value Chain Heatmap
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.5, 0.3, 12, 0.6, '투자 시사점 — 밸류체인별 시나리오 수혜도', 28, True)
add_text(slide, 0.5, 0.9, 12, 0.4, 'TSMC·SK하이닉스 Bull/Base 최고 수혜, 장비는 Bear에서도 방어적', 14, False, TEXT_G)

fig, ax = plt.subplots(figsize=(12, 4.5))
companies_inv = ['TSMC', 'SK하이닉스', 'NVIDIA', '삼성전자', 'ASML/장비', 'Shin-Etsu\n/SUMCO', '패키징']
bull_scores = [5, 5, 5, 4, 5, 4, 5]
base_scores = [5, 5, 4, 3, 4, 3, 4]
bear_scores = [3, 2, 2, 2, 3, 2, 2]

x_inv = np.arange(len(companies_inv))
w_inv = 0.25
ax.bar(x_inv - w_inv, bull_scores, w_inv, label='Bull', color=BULL_C, edgecolor='none')
ax.bar(x_inv, base_scores, w_inv, label='Base', color=BASE_C, edgecolor='none')
ax.bar(x_inv + w_inv, bear_scores, w_inv, label='Bear', color=BEAR_C, edgecolor='none')
ax.set_xticks(x_inv)
ax.set_xticklabels(companies_inv, fontsize=11)
ax.set_ylabel('수혜도 (★)', fontsize=12)
ax.set_ylim(0, 6)
ax.legend(fontsize=11)
ax.grid(axis='y', alpha=0.3)
ax.set_axisbelow(True)

# Add EPS notes
eps_notes = ['가동률1%p\n=EPS2%', 'HBM10%\n=영익15-20%', 'capex10%↓\n→EPS-15%', 'HBM(↑)\n파운드리(↓)', 'Backlog\n2-3년', 'LTA\n1-2년', 'CoWoS\n직결']
for i, note in enumerate(eps_notes):
    ax.text(i, 0.3, note, ha='center', va='bottom', fontsize=7, color=TEXT_G, style='italic')

plt.tight_layout()
add_chart_image(slide, fig, 0.5, 1.4, 12.3, 5.5)

# ══════════════════════════════════════════════
# SLIDE 12: Bear Reduction Order
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.5, 0.3, 12, 0.6, 'Bear 전환 시 포지션 축소 — 순서와 타이밍', 28, True)

fig, ax = plt.subplots(figsize=(12, 4.5))
# Timeline cascade
stages = [
    ('T+0', 'NVIDIA\n50%+ 축소', ACCENT_RED, 1.0),
    ('T+1-2Q', 'SK하이닉스\n40% 축소', ACCENT_RED, 0.8),
    ('T+2-3Q', '패키징\n30% 축소', ACCENT_YELLOW, 0.6),
    ('T+2-3Q', 'TSMC\n10-20%만', ACCENT_BLUE, 0.4),
    ('T+4-6Q', 'ASML/장비\n유지', ACCENT_GREEN, 0.3),
    ('T+4-6Q', '웨이퍼\n유지', ACCENT_GREEN, 0.2),
]
for i, (time, label, color, intensity) in enumerate(stages):
    ax.barh(5-i, 2+i*0.5, height=0.7, left=i*1.5, color=color, alpha=min(intensity+0.3, 1.0), edgecolor=BG, linewidth=1)
    ax.text(i*1.5 + (2+i*0.5)/2, 5-i, label, ha='center', va='center', fontsize=10, fontweight='bold', color='white')
    ax.text(i*1.5 - 0.1, 5-i, time, ha='right', va='center', fontsize=9, color=TEXT_G)

ax.set_xlim(-2, 12)
ax.set_ylim(-0.5, 6.5)
ax.axis('off')
ax.set_title('시그널 감지 → 실적 반영 타임라인', fontsize=13, pad=15, color=TEXT_W)

# 4 principles
principles = '4원칙: ① capex 민감도  ② 반영 속도  ③ 하방 보호  ④ 멀티플 리스크'
ax.text(5, -0.3, principles, ha='center', fontsize=11, color=ACCENT_YELLOW, fontweight='bold')
plt.tight_layout()
add_chart_image(slide, fig, 0.5, 1.4, 12.3, 5.5)

# ══════════════════════════════════════════════
# SLIDE 13: Monitoring Dashboard
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.5, 0.3, 12, 0.6, '모니터링 대시보드 — Trigger Points', 28, True)
add_text(slide, 0.5, 0.9, 12, 0.4, 'Bull 3개+ 충족 → 비중확대  |  Bear 2개+ 충족 → 축소실행', 14, False, TEXT_G)

fig, ax = plt.subplots(figsize=(12, 4.5))
indicators = [
    'TSMC AI/HPC 비중',
    '빅테크 capex',
    'HBM 계약가',
    'CoWoS 리드타임',
    'NVIDIA DC 매출',
    'DRAM DOI',
    '300mm 가동률',
]
current = [55, 280, 27.5, 18, 35, 4.5, 88]
bull_thresh = [60, 320, 35, 24, 45, 3, 93]
bear_thresh = [48, 240, 18, 10, 25, 7, 80]
# Normalize to 0-100 scale for visualization
def normalize(val, bear_t, bull_t):
    range_v = bull_t - bear_t
    if range_v == 0: return 50
    return max(0, min(100, (val - bear_t) / range_v * 100))

norm_current = [normalize(c, be, bu) for c, be, bu in zip(current, bear_thresh, bull_thresh)]

y_ind = np.arange(len(indicators))
# Background bar
ax.barh(y_ind, [100]*len(indicators), height=0.5, color='#222', edgecolor='none')
# Current position
for i, (nc, c) in enumerate(zip(norm_current, current)):
    color = ACCENT_GREEN if nc > 60 else ACCENT_YELLOW if nc > 30 else ACCENT_RED
    ax.barh(i, nc, height=0.5, color=color, alpha=0.7, edgecolor='none')
    ax.plot(nc, i, 'o', color='white', markersize=8, zorder=5)

# Bear/Bull lines
ax.axvline(0, color=ACCENT_RED, linewidth=2, linestyle='--', alpha=0.5, label='Bear 전환')
ax.axvline(100, color=ACCENT_GREEN, linewidth=2, linestyle='--', alpha=0.5, label='Bull 진입')

ax.set_yticks(y_ind)
ax.set_yticklabels(indicators, fontsize=11)
ax.set_xlim(-5, 105)
ax.set_xticks([0, 50, 100])
ax.set_xticklabels(['Bear', '현재', 'Bull'], fontsize=11)
ax.legend(fontsize=9, loc='lower right')
ax.invert_yaxis()
ax.set_title('핵심 지표 현재 위치 (Bear ↔ Bull 스펙트럼)', fontsize=13, pad=10)
plt.tight_layout()
add_chart_image(slide, fig, 0.5, 1.4, 12.3, 5.5)

# ══════════════════════════════════════════════
# SLIDE 14: Conclusion
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.5, 0.3, 12, 0.6, '결론 — 5가지 핵심', 28, True)

conclusions = [
    ('1', '병목 우선순위', 'CoWoS(26Q3 해소) → HBM(27까지 만성) → 원재료(Bull만)', ACCENT_BLUE),
    ('2', '최대 변수', 'Capex ±10%p → 웨이퍼 수요 16–22% 스윙', ACCENT_GREEN),
    ('3', '경쟁 구도', 'Intel/삼성 성공해도 총수요 불변 — TSMC 최상위 독점', ACCENT_PURPLE),
    ('4', '최대 리스크', 'AI버블+긴축 동시(8%) → 반도체 매출 –$150B', ACCENT_RED),
    ('5', '투자 핵심', 'Bear시 NVIDIA 1순위 축소(1–2Q), 장비/원재료 보유(4–6Q)', ACCENT_YELLOW),
]

for i, (num, title, desc, color) in enumerate(conclusions):
    y = 1.2 + i * 1.15
    add_shape_box(slide, 0.8, y, 11.7, 0.95, CARD_BG)
    # Number circle
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), Inches(y + 0.15), Inches(0.65), Inches(0.65))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.line.fill.background()
    tf = shape.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = rgb(TEXT_W)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    add_text(slide, 2.0, y + 0.1, 3, 0.4, title, 18, True, color)
    add_text(slide, 2.0, y + 0.5, 10, 0.4, desc, 14, False, TEXT_G)

add_text(slide, 0.5, 6.8, 12, 0.4, 'Semiconductor Research Team  |  2026.02.28  |  Confidential', 12, False, TEXT_G, PP_ALIGN.CENTER)

# ══════════════════════════════════════════════
# Save
# ══════════════════════════════════════════════
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'presentation_v2.pptx')
prs.save(output_path)
print(f'✅ Saved: {output_path}')
print(f'   Slides: {len(prs.slides)}')
